from pptx import Presentation


def load_ppt_text(file_path: str) -> str:
    prs = Presentation(file_path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text.strip()
